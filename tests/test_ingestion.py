"""Tests for document ingestion: parsers, chunker, and ingest registry."""

import os
import sys
import tempfile

import pytest


# ---------------------------------------------------------------------------
# Table handler
# ---------------------------------------------------------------------------

class TestTableHandler:
    """Tests for table_to_markdown conversion."""

    def test_normal_table(self):
        from src.ingestion.table_handler import table_to_markdown

        data = [
            ["Name", "Age", "City"],
            ["Alice", "30", "New York"],
            ["Bob", "25", "San Francisco"],
        ]
        md = table_to_markdown(data)
        assert "| Name | Age | City |" in md
        assert "| --- | --- | --- |" in md
        assert "| Alice | 30 | New York |" in md

    def test_empty_table(self):
        from src.ingestion.table_handler import table_to_markdown

        assert table_to_markdown([]) == ""
        assert table_to_markdown([[]]) == ""

    def test_ragged_table(self):
        from src.ingestion.table_handler import table_to_markdown

        ragged = [["A", "B"], ["1"]]
        md = table_to_markdown(ragged)
        assert "| 1 |  |" in md


# ---------------------------------------------------------------------------
# PDF parser
# ---------------------------------------------------------------------------

class TestPDFParser:
    """Tests for PDF parsing via pymupdf4llm."""

    def test_parse_pdf(self, tmp_path):
        import pymupdf
        from src.ingestion.pdf_parser import parse_pdf

        test_pdf = str(tmp_path / "test.pdf")
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text(
            (72, 72),
            "RAG System Architecture\n\n"
            "Retrieval-Augmented Generation combines retrieval with generation.",
            fontsize=12,
        )
        page2 = doc.new_page()
        page2.insert_text((72, 72), "Components of RAG\n\n1. Ingestion\n2. Retrieval", fontsize=12)
        doc.save(test_pdf)
        doc.close()

        documents = parse_pdf(test_pdf)
        assert len(documents) >= 2
        for d in documents:
            assert d["metadata"]["format"] == "pdf"
            assert "page_number" in d["metadata"]


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------

class TestPPTXParser:
    """Tests for PowerPoint parsing via python-pptx."""

    def test_parse_pptx(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        from src.ingestion.pptx_parser import parse_pptx

        test_pptx = str(tmp_path / "test.pptx")
        prs = Presentation()

        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        slide1.shapes.title.text = "Introduction to RAG"
        slide1.placeholders[1].text = "RAG enhances LLM responses."

        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = "Components"
        table_shape = slide2.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "Component"
        table.cell(0, 1).text = "Tool"
        table.cell(1, 0).text = "LLM"
        table.cell(1, 1).text = "Ollama"

        prs.save(test_pptx)

        documents = parse_pptx(test_pptx)
        assert len(documents) >= 2
        for d in documents:
            assert d["metadata"]["format"] == "pptx"


# ---------------------------------------------------------------------------
# OneNote parser
# ---------------------------------------------------------------------------

class TestOneNoteParser:
    """Tests for OneNote HTML export parsing."""

    def test_basic_html(self):
        from src.ingestion.onenote_parser import parse_onenote

        html = """<!DOCTYPE html>
<html><head><title>My OneNote Page</title></head>
<body>
<h1>Project Notes</h1>
<div><p>This is the first paragraph of my notes about the RAG system.</p>
<p>It contains important implementation details for the project.</p></div>
<h2>Architecture</h2>
<div><p>The system uses a modular architecture with separate components.</p></div>
<table><tr><th>Component</th><th>Tech</th></tr>
<tr><td>LLM</td><td>Ollama</td></tr></table>
</body></html>"""

        with tempfile.NamedTemporaryFile(mode="w", suffix=".htm", delete=False, encoding="utf-8") as f:
            f.write(html)
            tmp_path = f.name

        try:
            docs = parse_onenote(tmp_path)
            assert len(docs) > 0
            assert any(d["metadata"].get("page_title") == "My OneNote Page" for d in docs)
            assert any(d["metadata"].get("content_type") == "table" for d in docs)
            for d in docs:
                assert d["metadata"]["format"] == "onenote_html"
        finally:
            os.unlink(tmp_path)

    def test_empty_file(self):
        from src.ingestion.onenote_parser import parse_onenote

        with tempfile.NamedTemporaryFile(mode="w", suffix=".htm", delete=False, encoding="utf-8") as f:
            f.write("")
            tmp_path = f.name
        try:
            assert parse_onenote(tmp_path) == []
        finally:
            os.unlink(tmp_path)

    def test_chinese_content(self):
        from src.ingestion.onenote_parser import parse_onenote

        html = """<!DOCTYPE html>
<html><head><title>会议笔记</title></head>
<body><h1>项目讨论</h1>
<div><p>本次会议讨论了RAG系统的架构设计。</p>
<p>主要内容包括文档解析、向量检索和答案生成三个模块。</p></div>
</body></html>"""

        with tempfile.NamedTemporaryFile(mode="w", suffix=".html", delete=False, encoding="utf-8") as f:
            f.write(html)
            tmp_path = f.name
        try:
            docs = parse_onenote(tmp_path)
            assert len(docs) > 0
            assert any(d["metadata"].get("page_title") == "会议笔记" for d in docs)
        finally:
            os.unlink(tmp_path)

    def test_clean_text(self):
        from src.ingestion.onenote_parser import _clean_text

        assert _clean_text("line1\n\n\n\n\nline2") == "line1\n\nline2"
        assert _clean_text("  hello  \n  world  ") == "hello\nworld"

    def test_extract_tables(self):
        from bs4 import BeautifulSoup
        from src.ingestion.onenote_parser import _extract_tables

        html = """<table><tr><th>Name</th><th>Value</th></tr>
        <tr><td>Alpha</td><td>100</td></tr></table>"""
        soup = BeautifulSoup(html, "lxml")
        tables = _extract_tables(soup)
        assert len(tables) == 1
        assert "| Name | Value |" in tables[0]
        assert "| Alpha | 100 |" in tables[0]


# ---------------------------------------------------------------------------
# Chunker
# ---------------------------------------------------------------------------

class TestChunker:
    """Tests for document chunking."""

    def test_chunk_documents(self):
        from src.config import load_config
        from src.ingestion.chunker import DocumentChunker

        config = load_config()
        chunker = DocumentChunker(config)
        sample_docs = [
            {
                "content": "This is a sample page with enough text to test chunking. " * 20,
                "metadata": {"file_name": "test.pdf", "format": "pdf", "page_number": 1},
            },
        ]
        chunks = chunker.chunk_documents(sample_docs)
        assert len(chunks) >= 1
        for c in chunks:
            assert "chunk_id" in c
            assert "content" in c


# ---------------------------------------------------------------------------
# Ingest registry
# ---------------------------------------------------------------------------

class TestIngestRegistry:
    """Tests for the parser registry in ingest.py."""

    def test_all_parsers_registered(self):
        from src.ingest import PARSERS

        expected = {".txt", ".md", ".pdf", ".pptx", ".htm", ".html", ".ipynb"}
        actual = set(PARSERS.keys())
        assert actual == expected, f"Expected {expected}, got {actual}"
