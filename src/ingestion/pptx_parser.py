"""PowerPoint (.pptx) parser using python-pptx.

Extracts text, tables (as Markdown), and optionally describes images
via multimodal LLM. Each slide becomes one or more document dicts.
"""

import os
import time
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.logging.logger import get_logger

logger = get_logger("ingestion.pptx_parser")


def parse_pptx(
    file_path: str,
    extract_images: bool = False,
    config: Optional[dict] = None,
) -> List[Dict]:
    """Parse a PowerPoint file into document dicts with metadata.

    Processes each slide: extracts text from text frames, converts
    tables to Markdown, and optionally describes images via LLM.

    Args:
        file_path: Absolute path to the .pptx file.
        extract_images: If True, extract images and generate descriptions.
        config: Optional config dict (used for image LLM settings).

    Returns:
        List of dicts, each with keys: 'content', 'metadata'.
    """
    start = time.perf_counter()
    file_path = str(Path(file_path).resolve())

    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        return []

    file_name = os.path.basename(file_path)

    try:
        prs = Presentation(file_path)
    except Exception as e:
        logger.error(f"PPTX parsing failed for {file_name}: {e}")
        return []

    documents = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts = []
        slide_tables = []
        slide_images = []

        for shape in slide.shapes:
            # Text frames (titles, body text, etc.)
            if shape.has_text_frame:
                text = _extract_text_frame(shape.text_frame)
                if text:
                    slide_texts.append(text)

            # Tables
            if shape.has_table:
                md_table = _extract_table(shape.table)
                if md_table:
                    slide_tables.append(md_table)

            # Images
            if extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_desc = _extract_image(shape, config)
                if img_desc:
                    slide_images.append(img_desc)

        # Combine slide text content
        slide_content_parts = []

        if slide_texts:
            slide_content_parts.append("\n\n".join(slide_texts))

        if slide_tables:
            for i, table_md in enumerate(slide_tables, 1):
                slide_content_parts.append(f"[Table {i}]\n{table_md}")

        # Create text document for this slide
        if slide_content_parts:
            content = "\n\n".join(slide_content_parts)
            metadata = {
                "source_file": file_path,
                "file_name": file_name,
                "content_type": "text",
                "format": "pptx",
                "slide_number": slide_idx,
            }
            documents.append({"content": content, "metadata": metadata})

        # Create separate documents for table summaries
        for i, table_md in enumerate(slide_tables):
            summary = _summarize_table(table_md, config)
            if summary:
                metadata = {
                    "source_file": file_path,
                    "file_name": file_name,
                    "content_type": "table_summary",
                    "format": "pptx",
                    "slide_number": slide_idx,
                    "table_index": i,
                }
                documents.append({
                    "content": f"[Table {i + 1} summary, slide {slide_idx}]\n{summary}",
                    "metadata": metadata,
                })

        # Create separate documents for image descriptions
        for img_idx, img_desc in enumerate(slide_images):
            metadata = {
                "source_file": file_path,
                "file_name": file_name,
                "content_type": "image_description",
                "format": "pptx",
                "slide_number": slide_idx,
                "image_index": img_idx,
            }
            documents.append({
                "content": f"[Image on slide {slide_idx}]\n{img_desc}",
                "metadata": metadata,
            })

    elapsed_ms = int((time.perf_counter() - start) * 1000)
    logger.info(
        f"Parsed PPTX: {file_name} "
        f"({len(documents)} docs from {len(prs.slides)} slides, {elapsed_ms}ms)"
    )

    return documents


def _extract_text_frame(text_frame) -> str:
    """Extract text from a PowerPoint text frame.

    Args:
        text_frame: A pptx TextFrame object.

    Returns:
        Combined text from all paragraphs, stripped of excess whitespace.
    """
    paragraphs = []
    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _extract_table(table) -> str:
    """Extract a PowerPoint table and convert to Markdown.

    Args:
        table: A pptx Table object.

    Returns:
        Markdown-formatted table string.
    """
    # Lazy import to avoid circular dependency
    from src.ingestion.table_handler import table_to_markdown
    table_data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        table_data.append(row_data)

    if not table_data:
        return ""

    return table_to_markdown(table_data)


def _summarize_table(
    markdown_table: str,
    config: Optional[dict] = None,
) -> Optional[str]:
    """Generate a natural language summary of a table.

    Args:
        markdown_table: Markdown-formatted table string.
        config: Optional config dict for LLM settings.

    Returns:
        Summary string, or None on failure.
    """
    from src.ingestion.table_handler import summarize_table
    return summarize_table(markdown_table, config=config)


def _extract_image(shape, config: Optional[dict] = None) -> Optional[str]:
    """Extract and describe an image from a PowerPoint shape.

    Args:
        shape: A pptx Picture shape object.
        config: Optional config dict for LLM settings.

    Returns:
        Text description of the image, or None on failure.
    """
    # Lazy import to avoid circular dependency
    from src.ingestion.image_handler import describe_image_bytes

    try:
        image = shape.image
        image_bytes = image.blob
        content_type = image.content_type or "image/png"

        description = describe_image_bytes(
            image_bytes, mime_type=content_type, config=config
        )
        return description

    except Exception as e:
        logger.debug(f"Failed to extract image from shape: {e}")
        return None
